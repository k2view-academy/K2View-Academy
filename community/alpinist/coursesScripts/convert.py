"""Convert .docx and .pptx files in this folder to text-only markdown in ./md/.

Designed for vector-DB ingestion: clean structure (headings, lists, tables),
no images, no embedded media. Slide decks become one .md per deck with an H2
per slide; speaker notes are appended.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation


HERE = Path(__file__).parent
OUT = HERE / "md"
OUT.mkdir(exist_ok=True)


# ---------- docx ----------

def _para_md(para) -> str:
    text = "".join(run.text for run in para.runs).strip()
    if not text:
        return ""
    style = (para.style.name or "").lower() if para.style else ""
    if style.startswith("heading"):
        m = re.search(r"(\d+)", style)
        level = int(m.group(1)) if m else 1
        level = max(1, min(level, 6))
        return f"{'#' * level} {text}"
    if style.startswith("title"):
        return f"# {text}"
    if "list" in style or "bullet" in style:
        return f"- {text}"
    return text


def _table_md(table) -> str:
    rows = []
    for row in table.rows:
        cells = [re.sub(r"\s+", " ", cell.text).strip() for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    head = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    out = ["| " + " | ".join(head) + " |",
           "| " + " | ".join(["---"] * width) + " |"]
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _iter_block_items(parent):
    """Yield paragraphs and tables in document order from a docx body."""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = parent.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def convert_docx(path: Path) -> str:
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(path))
    parts: list[str] = [f"# {path.stem}", ""]
    for block in _iter_block_items(doc):
        if isinstance(block, Paragraph):
            md = _para_md(block)
            if md:
                parts.append(md)
                parts.append("")
        elif isinstance(block, Table):
            md = _table_md(block)
            if md:
                parts.append(md)
                parts.append("")
    return "\n".join(parts).rstrip() + "\n"


# ---------- pptx ----------

def _shape_text(shape) -> list[str]:
    lines: list[str] = []
    if not shape.has_text_frame:
        # Group shapes can contain children
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for sub in shape.shapes:
                lines.extend(_shape_text(sub))
        return lines
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        if para.level and para.level > 0:
            lines.append("  " * para.level + f"- {text}")
        else:
            lines.append(text)
    return lines


def _slide_md(slide, idx: int) -> str:
    title = ""
    body_lines: list[str] = []
    table_blocks: list[str] = []

    if slide.shapes.title and slide.shapes.title.text.strip():
        title = slide.shapes.title.text.strip()

    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([re.sub(r"\s+", " ",
                                    "".join(c.text for c in row.cells if c is not None)).strip()
                             for row in [row]
                             for c in [row]] if False else
                            [re.sub(r"\s+", " ", cell.text).strip() for cell in row.cells])
            if rows:
                width = max(len(r) for r in rows)
                rows = [r + [""] * (width - len(r)) for r in rows]
                tbl = ["| " + " | ".join(rows[0]) + " |",
                       "| " + " | ".join(["---"] * width) + " |"]
                for r in rows[1:]:
                    tbl.append("| " + " | ".join(r) + " |")
                table_blocks.append("\n".join(tbl))
            continue
        for line in _shape_text(shape):
            body_lines.append(line)

    parts = [f"## Slide {idx}" + (f": {title}" if title else "")]
    if body_lines:
        parts.append("")
        # Heuristic: turn plain lines into bullets if there are multiple
        if len(body_lines) > 1 and not any(l.lstrip().startswith("-") for l in body_lines):
            parts.extend(f"- {l}" for l in body_lines)
        else:
            parts.extend(body_lines)
    for tbl in table_blocks:
        parts.append("")
        parts.append(tbl)

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append("")
            parts.append("**Notes:**")
            parts.append("")
            parts.append(notes)

    return "\n".join(parts).rstrip()


def convert_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts = [f"# {path.stem}", ""]
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(_slide_md(slide, i))
        parts.append("")
    return "\n".join(parts).rstrip() + "\n"


# ---------- driver ----------

def main() -> int:
    files = sorted([p for p in HERE.iterdir()
                    if p.suffix.lower() in (".docx", ".pptx") and not p.name.startswith("~$")])
    if not files:
        print("No .docx/.pptx files found.")
        return 1

    failures: list[tuple[Path, str]] = []
    for src in files:
        out = OUT / (src.stem + ".md")
        try:
            if src.suffix.lower() == ".docx":
                md = convert_docx(src)
            else:
                md = convert_pptx(src)
            out.write_text(md, encoding="utf-8")
            print(f"OK   {src.name} -> md/{out.name} ({len(md):,} chars)")
        except Exception as e:
            failures.append((src, repr(e)))
            print(f"FAIL {src.name}: {e!r}")

    if failures:
        print(f"\n{len(failures)} file(s) failed.")
        return 2
    print(f"\nConverted {len(files)} file(s) into {OUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
